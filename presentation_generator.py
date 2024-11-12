from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from typing import Optional, List, Dict, Union
import pandas as pd
import tempfile
import logging
import os

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

def add_chart_to_slide(slide, chart_data: dict, left: float, top: float, width: float, height: float):
    """Add a chart to the slide using native PowerPoint charts."""
    try:
        # Create chart data
        chart_data_obj = CategoryChartData()
        chart_data_obj.categories = chart_data['categories']
        
        # Add series
        for series_name, values in chart_data['series'].items():
            chart_data_obj.add_series(series_name, values)
        
        # Map visualization types to PowerPoint chart types
        chart_type_map = {
            'bar': XL_CHART_TYPE.COLUMN_CLUSTERED,
            'line': XL_CHART_TYPE.LINE,
            'pie': XL_CHART_TYPE.PIE,
            'scatter': XL_CHART_TYPE.XY_SCATTER,
        }
        
        # Add chart to slide
        chart_type = chart_type_map.get(chart_data['type'], XL_CHART_TYPE.COLUMN_CLUSTERED)
        chart = slide.shapes.add_chart(
            chart_type,
            Inches(left), Inches(top),
            Inches(width), Inches(height),
            chart_data_obj
        ).chart
        
        # Set title and labels
        chart.has_title = True
        chart.chart_title.text_frame.text = chart_data['title']
        
        if hasattr(chart, 'value_axis'):
            chart.value_axis.has_major_gridlines = True
            chart.value_axis.has_minor_gridlines = False
            chart.value_axis.title.text_frame.text = chart_data['y_label']
            
        if hasattr(chart, 'category_axis'):
            chart.category_axis.has_major_gridlines = False
            chart.category_axis.has_minor_gridlines = False
            chart.category_axis.title.text_frame.text = chart_data['x_label']
        
        return chart
    except Exception as e:
        logger.error(f"Error adding chart to slide: {str(e)}")
        raise Exception(f"Error adding chart to slide: {str(e)}")

def create_presentation(data: pd.DataFrame, summaries: List[str], visualizations: Optional[List[Dict]] = None) -> str:
    """
    Create a PowerPoint presentation with the processed data, AI summaries, and visualizations.
    """
    prs = Presentation()
    
    # Title Slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    if len(title_slide.shapes.title.text_frame.paragraphs) > 0:
        title_slide.shapes.title.text_frame.text = "Data Analysis Report"
    if len(title_slide.placeholders) > 1:
        subtitle = title_slide.placeholders[1]
        if hasattr(subtitle, "text_frame"):
            subtitle.text_frame.text = "Generated with AI Insights"

    # Summary Slide
    summary_slide = prs.slides.add_slide(prs.slide_layouts[1])
    if hasattr(summary_slide.shapes.title, "text_frame"):
        summary_slide.shapes.title.text_frame.text = "Key Insights"
    
    # Add summary content
    content_placeholder = None
    for shape in summary_slide.placeholders:
        if shape.placeholder_format.idx == 1:  # Content placeholder
            content_placeholder = shape
            break
    
    if content_placeholder and hasattr(content_placeholder, "text_frame"):
        content_placeholder.text_frame.text = "\n".join([f"• {summary}" for summary in summaries])

    # Visualization Slides
    if visualizations:
        for idx, viz_data in enumerate(visualizations):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            if hasattr(slide.shapes.title, "text_frame"):
                slide.shapes.title.text_frame.text = viz_data.get('title', f"Visualization {idx + 1}")
            
            # Add chart to slide using native PowerPoint charts
            try:
                add_chart_to_slide(slide, viz_data, left=1, top=2, width=8, height=5)
            except Exception as e:
                logger.error(f"Error adding visualization {idx + 1} to slide: {str(e)}")
                # Add error message to slide if chart creation fails
                txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
                tf = txBox.text_frame
                tf.text = f"Error adding visualization: {str(e)}"

    # Data Slides
    for i in range(0, len(data), 5):  # 5 rows per slide
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        
        # Add title
        if hasattr(slide.shapes.title, "text_frame"):
            slide.shapes.title.text_frame.text = "Data Analysis"
        
        # Add table
        rows = min(6, len(data) - i + 1)  # +1 for header
        cols = len(data.columns)
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(0.8 * rows)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Add headers
        for j, col in enumerate(data.columns):
            cell = table.cell(0, j)
            cell.text = str(col)
            
        # Add data
        for row in range(rows-1):
            if i + row < len(data):
                for col in range(cols):
                    cell = table.cell(row+1, col)
                    cell.text = str(data.iloc[i + row, col])

    # Save presentation
    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
    prs.save(temp_file.name)
    return temp_file.name
